"""
trusscore_deck.py — canonical Trusscore PowerPoint builder.

Import this when building any Trusscore deck. It enforces the brand spec in code:
type ladder (auto-tier), tight bullet spacing, two-level truss bullets (solid L1 /
outline L2), no terminal periods on bullets, the locked palette, and the master template.

Usage:
    import trusscore_deck as td
    prs = td.open_deck()
    td.cover(prs, "AI Output Brand Alignment Playbook", "How Trusscore uses Claude...", yellow="Playbook")
    td.agenda(prs, ["Quick start", "What the system is", "How to use it"])
    td.section(prs, "What We Built", "The brand now shapes what Claude builds...")
    td.content(prs, "The Three Pieces", [("Brand guidelines","The rulebook: one source of truth"), ...])
    td.cards(prs, "Quick Start", [("Ask normally","Plain language","Request work like a teammate"), ...], slate=1)
    td.stats(prs, "What This Means", [("1","source of truth","one place"), ...], slate=1)
    td.flow(prs, "How They Work Together", [("Brand guidelines","The rulebook"), ...], callout="One source...")
    td.close(prs, "On Brand by Default", "Just ask, then de-slop it")
    td.save(prs, "/path/out.pptx")

Sizing rule (Section 14b): body auto-tiers 28/24 -> 24/20 -> 20/16 (16pt hard floor);
if a slide won't fit at the floor, pick_tier returns the floor and you should SPLIT the
slide or cut copy rather than shipping crowded text.
"""
import os
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

# ---- palette (locked) ----
SLATE=RGBColor(0x3A,0x4B,0x5C); TEXT=RGBColor(0x11,0x1A,0x22); GREY=RGBColor(0x6A,0x7A,0x88)
LIGHT=RGBColor(0xF5,0xF6,0xF7); BORDER=RGBColor(0xE1,0xE5,0xE8); YELLOW=RGBColor(0xFE,0xB1,0x00)
WHITE=RGBColor(0xFF,0xFF,0xFF); LGREY=RGBColor(0xDD,0xE3,0xE8); TINT=RGBColor(0xFF,0xF8,0xEC)
GREEN=RGBColor(0x6B,0xA5,0x43); RED=RGBColor(0xDC,0x35,0x45)   # functional status only

TEMPLATE=os.path.join(os.path.dirname(__file__),"..","assets","Trusscore_Powerpoint_Template.pptx")
X3=[0.74,4.79,8.84]; CY=2.45; CW=3.75; CH=3.15          # 3-up card geometry
LADDER=[(28,24),(24,20),(20,16)]                         # top/sub pt; 16 = hard floor

# ---------- core ----------
def open_deck(template=TEMPLATE):
    prs=Presentation(template)
    lst=prs.slides._sldIdLst
    for sid in list(lst):
        try: prs.part.drop_rel(sid.get(qn('r:id')))
        except Exception: pass
        lst.remove(sid)
    prs._td_layouts={l.name:l for l in prs.slide_layouts}
    return prs

def _add(prs,name): return prs.slides.add_slide(prs._td_layouts[name])
def save(prs,path): prs.save(path); return path

def _ctitle(slide,t):
    slide.shapes.title.text=t
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.name="Aptos Light"; r.font.size=Pt(48); r.font.bold=False; r.font.color.rgb=SLATE

def _setind(p,marL,indent):
    pPr=p._p.get_or_add_pPr(); pPr.set('marL',str(marL)); pPr.set('indent',str(indent))

def pick_tier(items):
    """Largest ladder step where all two-level items fit the content area; else the floor."""
    tuples=[it for it in items if isinstance(it,tuple)]
    if not tuples: return (28,24)
    usable,width,cf,lh=4.5,11.2,0.47,1.2   # inches; calibrated to renders
    for lab,sub in LADDER:
        total=0.0
        for _,desc in tuples:
            total += lab*lh/72 + (6+3)/72
            cpl=max(1,int(width/(sub*cf/72))); lines=max(1,-(-len(str(desc).strip())//cpl))
            total += lines*sub*lh/72 + 6/72
        if total<=usable: return (lab,sub)
    return LADDER[-1]

def _body(slide,idx,items):
    tf=slide.placeholders[idx].text_frame; tf.clear(); first=True
    lab_pt,sub_pt=pick_tier(items)
    for it in items:
        if isinstance(it,tuple):
            label=str(it[0]).strip().rstrip('.').strip(); desc=str(it[1]).strip().rstrip('.')
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            p.level=0; p.space_after=Pt(3); p.space_before=Pt(6); _setind(p,256032,-256032)
            r=p.add_run(); r.text=label; r.font.bold=True; r.font.name="Aptos"; r.font.size=Pt(lab_pt)
            p2=tf.add_paragraph(); p2.level=1; p2.space_after=Pt(6); p2.space_before=Pt(0); _setind(p2,512064,-256032)
            r2=p2.add_run(); r2.text=desc; r2.font.name="Aptos"; r2.font.size=Pt(sub_pt)
        else:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            p.level=0; _setind(p,256032,-256032); p.add_run().text=str(it).strip().rstrip('.')
    return lab_pt,sub_pt

def _rrect(slide,x,y,w,h,fill,line=None,radius=0.06):
    sp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sp.adjustments[0]=radius; sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line: sp.line.color.rgb=line; sp.line.width=Pt(0.75)
    else: sp.line.fill.background()
    sp.shadow.inherit=False
    return sp

def _accent(slide,cx,cy):
    a=slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,Inches(cx),Inches(cy),Inches(0.34),Inches(0.17))
    a.fill.solid(); a.fill.fore_color.rgb=YELLOW; a.line.fill.background(); a.shadow.inherit=False

def _cardtext(shape,rows,anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT):
    tf=shape.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.24); tf.margin_right=Inches(0.2); tf.margin_top=Inches(0.22)
    for i,(t,sz,col,b,sa) in enumerate(rows):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.space_after=Pt(sa)
        r=p.add_run(); r.text=t; r.font.name="Aptos"; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=col

# ---------- slide builders ----------
def cover(prs,title,subtitle="",yellow=""):
    s=_add(prs,'Title1'); tp=s.placeholders[0].text_frame; tp.clear(); p=tp.paragraphs[0]
    if yellow and yellow in title:
        pre,post=title.split(yellow,1)
        if pre: r=p.add_run(); r.text=pre; r.font.color.rgb=WHITE
        r=p.add_run(); r.text=yellow; r.font.color.rgb=YELLOW
        if post: r=p.add_run(); r.text=post; r.font.color.rgb=WHITE
    else:
        r=p.add_run(); r.text=title; r.font.color.rgb=WHITE
    if subtitle:
        for idx in (16,11):
            try: s.placeholders[idx].text_frame.text=subtitle; break
            except KeyError: pass
    return s

def agenda(prs,items):
    s=_add(prs,'Agenda'); _ctitle(s,"Agenda"); _body(s,1,items); return s

def section(prs,title,statement):
    s=_add(prs,'Section - General'); s.shapes.title.text=title; _body(s,14,[statement]); return s

def content(prs,title,items):
    """Title and Content. items = list of (label, desc) two-level bullets or plain strings. Auto-tiers."""
    s=_add(prs,'Title and Content'); _ctitle(s,title); _body(s,1,items); return s

def cards(prs,title,items,slate=None):
    """3-up card row. items = list of (name, role, desc). slate = index (0-2) of the one dark emphasis card."""
    s=_add(prs,'Title Only'); _ctitle(s,title)
    for i,(name,role,desc) in enumerate(items[:3]):
        is_slate=(i==slate); fill=SLATE if is_slate else LIGHT
        c=_rrect(s,X3[i],CY,CW,CH,fill,None if is_slate else BORDER,radius=0.06)
        nameC=WHITE if is_slate else SLATE; descC=LGREY if is_slate else GREY
        _cardtext(c,[(name,22,nameC,True,4),(role,16,YELLOW,True,10),(desc,16,descC,False,0)])
        if not is_slate: _accent(s,X3[i]+CW-0.5,CY+0.22)
    return s

def stats(prs,title,items,slate=None):
    """3-up stat callouts. items = list of (big_number, label, sub). slate = index of emphasis card."""
    s=_add(prs,'Title Only'); _ctitle(s,title)
    for i,(num,label,sub) in enumerate(items[:3]):
        is_slate=(i==slate); fill=SLATE if is_slate else LIGHT
        c=_rrect(s,X3[i],CY,CW,CH,fill,None if is_slate else BORDER,radius=0.06)
        numC=YELLOW if is_slate else SLATE; labC=WHITE if is_slate else SLATE; subC=LGREY if is_slate else GREY
        _cardtext(c,[(num,54,numC,True,2),(label,16,labC,True,6),(sub,16,subC,False,0)],
                  anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    return s

def flow(prs,title,nodes,callout=None):
    """Left-to-right flow of up to 3 nodes = (name, role), joined by yellow arrows, optional callout."""
    s=_add(prs,'Title Only'); _ctitle(s,title)
    FW,FH,FY=3.5,2.15,2.15; fx=[0.72,4.92,9.12]
    for i,(name,role) in enumerate(nodes[:3]):
        c=_rrect(s,fx[i],FY,FW,FH,WHITE,BORDER,radius=0.06)
        tf=c.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        for j,(t,sz,col,b) in enumerate([(name,20,SLATE,True),(role,16,YELLOW,True)]):
            p=tf.paragraphs[0] if j==0 else tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER
            r=p.add_run(); r.text=t; r.font.name="Aptos"; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=col
    for ax in (4.27,8.47)[:max(0,len(nodes)-1)]:
        a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(ax),Inches(2.95),Inches(0.6),Inches(0.45))
        a.fill.solid(); a.fill.fore_color.rgb=YELLOW; a.line.fill.background(); a.shadow.inherit=False
    if callout: _callout(s,callout)
    return s

def status_row(prs,title,items):
    """3-up verdict/status cards. items = list of (label, color, desc). color: 'green'|'yellow'|'red'|RGBColor."""
    s=_add(prs,'Title Only'); _ctitle(s,title)
    cmap={'green':GREEN,'yellow':YELLOW,'red':RED,'slate':SLATE}
    for i,(label,color,desc) in enumerate(items[:3]):
        col=cmap.get(color,color) if not isinstance(color,RGBColor) else color
        c=_rrect(s,X3[i],CY,CW,CH-0.4,LIGHT,BORDER,radius=0.06)
        _cardtext(c,[(label,22,col,True,8),(desc,16,GREY,False,0)],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    return s

def _callout(slide,text,y=4.75):
    c=_rrect(slide,0.72,y,11.9,0.95,TINT,YELLOW,radius=0.06)
    tf=c.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(text).strip().rstrip('.'); r.font.name="Aptos"; r.font.size=Pt(16); r.font.bold=True; r.font.color.rgb=SLATE

def callout(slide,text,y=4.75): _callout(slide,text,y)   # add a callout to an existing slide

def close(prs,headline,subline=""):
    s=_add(prs,'Ending1'); s.placeholders[0].text=headline
    for p in s.placeholders[0].text_frame.paragraphs:
        for r in p.runs: r.font.color.rgb=WHITE
    if subline:
        sub=s.shapes.add_textbox(Inches(0.72),Inches(4.55),Inches(9),Inches(0.6))
        p=sub.text_frame.paragraphs[0]; r=p.add_run(); r.text=subline
        r.font.name="Aptos"; r.font.size=Pt(20); r.font.bold=True; r.font.color.rgb=YELLOW
    return s
